import sys
import json
import requests
import tempfile
import os
import uuid
from pptx import Presentation

try:
    import comtypes.client
except ImportError:
    comtypes = None

def replace_in_shape(shape, replacements):
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            replace_in_shape(s, replacements)
            
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame"):
                    _replace_in_text_frame(cell.text_frame, replacements)
                    
    if hasattr(shape, "text_frame"):
        _replace_in_text_frame(shape.text_frame, replacements)

def _replace_in_text_frame(text_frame, replacements):
    for paragraph in text_frame.paragraphs:
        original_text = paragraph.text
        new_text = original_text
        for key, value in replacements.items():
            if key in new_text:
                new_text = new_text.replace(key, str(value) if value is not None else "")
                
        if new_text != original_text:
            if paragraph.runs:
                first_run = paragraph.runs[0]
                first_run.text = new_text
                for i in range(len(paragraph.runs) - 1, 0, -1):
                    p = paragraph._p
                    p.remove(paragraph.runs[i]._r)


def main():
    try:
        input_data = json.loads(sys.stdin.read())
        
        url = input_data.get("templateUrl")
        output_dir = input_data.get("outputDir", tempfile.gettempdir())
        replacements = input_data.get("replacements", {})
        
        if not url:
            raise ValueError("Template URL is missing")
        
        # Download
        response = requests.get(url)
        response.raise_for_status()
        
        file_id = str(uuid.uuid4())
        tmp_in_path = os.path.join(output_dir, f"{file_id}_in.pptx")
        
        with open(tmp_in_path, "wb") as f:
            f.write(response.content)
            
        prs = Presentation(tmp_in_path)
        
        # Replace placeholders
        for slide in prs.slides:
            for shape in slide.shapes:
                replace_in_shape(shape, replacements)
                            
        # Save generated pptx
        out_pptx_name = f"{file_id}.pptx"
        out_pptx_path = os.path.join(output_dir, out_pptx_name)
        prs.save(out_pptx_path)
        sys.stderr.write(f"Generated PPTX path: {out_pptx_path}\n")
        
        # Export ALL slides as images using comtypes
        img_files = []

        if comtypes:
            powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
            pres_export = powerpoint.Presentations.Open(os.path.abspath(out_pptx_path), WithWindow=False)
            
            # ── Export every slide ──────────────────────────────────────────
            for i in range(1, pres_export.Slides.Count + 1):
                slide_img_name = f"{file_id}_slide{i}.jpg"
                slide_img_path = os.path.join(output_dir, slide_img_name)
                pres_export.Slides[i].Export(os.path.abspath(slide_img_path), "JPG")
                img_files.append(slide_img_name)
                sys.stderr.write(f"Exported slide {i}: {slide_img_path}\n")

            pres_export.Close()
            powerpoint.Quit()
        else:
            raise Exception("comtypes not available for image generation")
        
        # Cleanup temp input
        if os.path.exists(tmp_in_path):
            os.remove(tmp_in_path)
            
        print(json.dumps({
            "success":   True,
            "pptx_file": out_pptx_name,
            "img_files": img_files,                          # all slide images
            "img_file":  img_files[0] if img_files else ""  # backward compat
        }))
        
    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()