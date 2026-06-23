import sys
import json
import requests
import tempfile
import os
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

try:
    import comtypes.client
except ImportError:
    comtypes = None

# No font mapping needed — frontend now sends exact PowerPoint font names
# e.g. "Calibri", "Georgia", "Arial" — same name works in CSS and PPT.

# PowerPoint "SaveAs" format constant for PDF (ppSaveAsPDF).
PP_SAVE_AS_PDF = 32


def hex_to_rgb(hex_color):
    """Convert #RRGGBB to RGBColor, return None on failure."""
    try:
        h = hex_color.lstrip('#')
        if len(h) == 6:
            r, g, b = int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16)
            return RGBColor(r, g, b)
    except Exception:
        pass
    return None

def apply_style_to_run(run, style):
    """Apply font styling from the style dict to a single run."""
    if not style:
        return
    try:
        font = run.font

        if style.get('bold') is not None:
            font.bold = bool(style['bold'])

        if style.get('italic') is not None:
            font.italic = bool(style['italic'])

        if style.get('underline') is not None:
            font.underline = bool(style['underline'])

        if style.get('fontSize'):
            font.size = Pt(float(style['fontSize']))

        if style.get('color'):
            rgb = hex_to_rgb(style['color'])
            if rgb:
                font.color.rgb = rgb

        # fontFamily value is already the exact PPT font name — use directly, no mapping
        if style.get('fontFamily'):
            font.name = style['fontFamily']

    except Exception as e:
        sys.stderr.write(f"Style apply error: {e}\n")

def replace_in_shape(shape, replacements, styling):
    if hasattr(shape, "shapes"):
        for s in shape.shapes:
            replace_in_shape(s, replacements, styling)

    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame"):
                    _replace_in_text_frame(cell.text_frame, replacements, styling)

    if hasattr(shape, "text_frame"):
        _replace_in_text_frame(shape.text_frame, replacements, styling)

def _replace_in_text_frame(text_frame, replacements, styling):
    for paragraph in text_frame.paragraphs:
        # paragraph.text merges all runs — use it to detect matches even when
        # a placeholder like {{Lorem Ipsum 2}} is split across multiple runs.
        para_text   = paragraph.text
        matched_key = None
        for key in replacements:
            if key in para_text:
                matched_key = key
                break

        if not matched_key:
            continue

        new_value = str(replacements[matched_key]) if replacements[matched_key] is not None else ""
        style     = styling.get(matched_key, {})

        runs = paragraph.runs
        if not runs:
            continue

        # Merge ALL run texts into one string (handles split-run placeholders),
        # do the replacement, put the result in run[0], blank every other run.
        combined = "".join(r.text for r in runs)

        # If the placeholder isn't in the merged runs but IS in paragraph.text,
        # it may live inside XML elements python-pptx doesn't surface as runs
        # (e.g. field elements). Fall back to para_text in that case.
        source = combined if matched_key in combined else para_text
        if matched_key not in source:
            continue

        new_combined = source.replace(matched_key, new_value)

        runs[0].text = new_combined
        for r in runs[1:]:
            r.text = ""

        # Apply styling to every run in the paragraph
        for r in paragraph.runs:
            apply_style_to_run(r, style)


def export_pdf(pptx_path, output_dir):
    """Open the generated PPTX in PowerPoint via COM and export it as a PDF.
    Returns the absolute path to the produced PDF file."""
    if not comtypes:
        raise Exception("comtypes not available for PDF export")

    pdf_path = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf")

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    try:
        pres = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        try:
            pres.SaveAs(os.path.abspath(pdf_path), PP_SAVE_AS_PDF)
            sys.stderr.write(f"Exported PDF: {pdf_path}\n")
        finally:
            pres.Close()
    finally:
        powerpoint.Quit()

    return pdf_path


def export_image(pptx_path, output_dir):
    """Open the generated PPTX in PowerPoint via COM and export slide 1 as a JPG.
    Returns the absolute path to the produced JPG file."""
    if not comtypes:
        raise Exception("comtypes not available for image export")

    img_path = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_path))[0] + "_slide1.jpg")

    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    try:
        pres = powerpoint.Presentations.Open(os.path.abspath(pptx_path), WithWindow=False)
        try:
            # Only the first slide is needed for the "Image" download option.
            pres.Slides[1].Export(os.path.abspath(img_path), "JPG")
            sys.stderr.write(f"Exported image: {img_path}\n")
        finally:
            pres.Close()
    finally:
        powerpoint.Quit()

    return img_path


def main():
    try:
        input_data = json.loads(sys.stdin.read())

        url = input_data.get("templateUrl")
        if not url:
            raise ValueError("Template URL is missing")

        replacements = input_data.get("replacements", {})
        styling      = input_data.get("styling", {})
        # 'pptx' (default) | 'pdf' | 'image'
        out_format   = (input_data.get("format") or "pptx").lower()

        response = requests.get(url)
        response.raise_for_status()

        with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_in:
            tmp_in.write(response.content)
            tmp_in_path = tmp_in.name

        prs = Presentation(tmp_in_path)

        for slide in prs.slides:
            for shape in slide.shapes:
                replace_in_shape(shape, replacements, styling)

        tmp_out_path = tmp_in_path.replace(".pptx", "_out.pptx")
        prs.save(tmp_out_path)

        if os.path.exists(tmp_in_path):
            os.remove(tmp_in_path)

        final_path   = tmp_out_path
        final_format = "pptx"

        # ── Convert the finished, replacement-applied deck to PDF/Image ────
        # via PowerPoint COM, reusing the same approach as generate_preview.py.
        if out_format == "pdf":
            pdf_path = export_pdf(tmp_out_path, os.path.dirname(tmp_out_path))
            final_path   = pdf_path
            final_format = "pdf"
        elif out_format == "image":
            img_path = export_image(tmp_out_path, os.path.dirname(tmp_out_path))
            final_path   = img_path
            final_format = "image"

        print(json.dumps({
            "success":     True,
            "output_path": final_path,
            "format":      final_format,
            # Keep the intermediate pptx path so the caller can clean it up
            # too when a conversion produced a separate pdf/jpg file.
            "pptx_path":   tmp_out_path,
        }))

    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()