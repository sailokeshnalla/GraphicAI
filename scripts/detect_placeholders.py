import sys
import json
import requests
import tempfile
import os
import re
from xml.etree import ElementTree as ET

def extract_svg_placeholders(svg_content):
    """Extract {{placeholder}} patterns from SVG text elements"""
    placeholders_map = {}

    try:
        # Parse SVG
        root = ET.fromstring(svg_content)
        ns = {
            'svg': 'http://www.w3.org/2000/svg',
            'xlink': 'http://www.w3.org/1999/xlink'
        }

        # Get SVG dimensions
        viewBox = root.get('viewBox', '')
        width_attr = root.get('width', '1000')
        height_attr = root.get('height', '1000')

        svg_width = float(re.sub(r'[^\d.]', '', str(width_attr)) or 1000)
        svg_height = float(re.sub(r'[^\d.]', '', str(height_attr)) or 1000)

        if viewBox:
            parts = viewBox.strip().split()
            if len(parts) == 4:
                svg_width = float(parts[2])
                svg_height = float(parts[3])

        # Find all text/tspan elements (SVG text lives here)
        all_text_elements = root.iter()

        for elem in all_text_elements:
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag
            if tag not in ('text', 'tspan', 'flowRoot', 'flowPara'):
                continue

            # Collect full text including tspan children
            full_text = (elem.text or '')
            for child in elem.iter():
                if child is not elem:
                    full_text += (child.text or '') + (child.tail or '')

            matches = re.findall(r"\{\{.*?\}\}", full_text)
            if not matches:
                continue

            # Get position - walk up to find a parent with x/y if needed
            x = float(elem.get('x', 0) or 0)
            y = float(elem.get('y', 0) or 0)

            # Try transform translate if no direct x/y
            transform = elem.get('transform', '')
            tx, ty = 0.0, 0.0
            t_match = re.search(r'translate\(([-\d.]+)[,\s]+([-\d.]+)\)', transform)
            if t_match:
                tx = float(t_match.group(1))
                ty = float(t_match.group(2))

            x += tx
            y += ty

            # Font styling
            font_size = float(re.sub(r'[^\d.]', '', elem.get('font-size', '18')) or 18)
            font_color = elem.get('fill', '#1e293b')
            font_weight = elem.get('font-weight', 'normal')
            font_style = elem.get('font-style', 'normal')
            text_anchor = elem.get('text-anchor', 'start')

            align_map = {'middle': 'CENTER', 'end': 'RIGHT', 'start': 'LEFT'}
            align_val = align_map.get(text_anchor, 'LEFT')

            # Estimate width from font size
            estimated_width = max(len(full_text) * font_size * 0.6, font_size * 5)

            left_pct = (x / svg_width) * 100
            top_pct = (y / svg_height) * 100
            width_pct = min((estimated_width / svg_width) * 100, 50)
            height_pct = (font_size * 1.5 / svg_height) * 100

            for match in matches:
                if match not in placeholders_map:
                    placeholders_map[match] = {
                        "left": round(left_pct, 2),
                        "top": round(top_pct, 2),
                        "width": round(width_pct, 2),
                        "height": round(height_pct, 2),
                        "fontSize": font_size,
                        "align": align_val,
                        "color": font_color if font_color.startswith('#') else '#1e293b',
                        "isBold": font_weight in ('bold', '700', '800', '900'),
                        "isItalic": font_style == 'italic'
                    }

    except ET.ParseError as e:
        raise ValueError(f"Invalid SVG file: {e}")

    return placeholders_map, svg_width, svg_height


def extract_placeholders(shape, slide_width, slide_height, placeholders_map, scale_x=1.0, scale_y=1.0, offset_x=0.0, offset_y=0.0):
    if hasattr(shape, "shapes"):
        group_left_slide = (shape.left or 0) * scale_x + offset_x
        group_top_slide = (shape.top or 0) * scale_y + offset_y
        
        chOff_x, chOff_y = 0, 0
        chExt_cx, chExt_cy = 1, 1
        try:
            chOff = shape._element.xpath("./p:grpSpPr/a:xfrm/a:chOff")
            if chOff:
                chOff_x = int(chOff[0].get("x"))
                chOff_y = int(chOff[0].get("y"))
            else:
                chOff_x = shape.left or 0
                chOff_y = shape.top or 0
                
            chExt = shape._element.xpath("./p:grpSpPr/a:xfrm/a:chExt")
            if chExt:
                chExt_cx = int(chExt[0].get("cx"))
                chExt_cy = int(chExt[0].get("cy"))
            else:
                chExt_cx = shape.width or 1
                chExt_cy = shape.height or 1
        except Exception:
            chOff_x = shape.left or 0
            chOff_y = shape.top or 0
            chExt_cx = shape.width or 1
            chExt_cy = shape.height or 1
            
        new_scale_x = ((shape.width or 1) / (chExt_cx or 1)) * scale_x
        new_scale_y = ((shape.height or 1) / (chExt_cy or 1)) * scale_y
        new_offset_x = group_left_slide - chOff_x * new_scale_x
        new_offset_y = group_top_slide - chOff_y * new_scale_y
        
        for s in shape.shapes:
            extract_placeholders(s, slide_width, slide_height, placeholders_map, new_scale_x, new_scale_y, new_offset_x, new_offset_y)
            
    if hasattr(shape, "table"):
        for row in shape.table.rows:
            for cell in row.cells:
                if hasattr(cell, "text_frame"):
                    _extract_from_text_frame(cell.text_frame, shape, slide_width, slide_height, placeholders_map, scale_x, scale_y, offset_x, offset_y)

    if hasattr(shape, "text_frame"):
        _extract_from_text_frame(shape.text_frame, shape, slide_width, slide_height, placeholders_map, scale_x, scale_y, offset_x, offset_y)

def _extract_from_text_frame(text_frame, shape, slide_width, slide_height, placeholders_map, scale_x, scale_y, offset_x, offset_y):
    for paragraph in text_frame.paragraphs:
        # Use paragraph.text (merged across all runs) to detect placeholders.
        # A single placeholder like {{Lorem Ipsum 2}} can be split across
        # multiple runs in the XML, so searching run-by-run misses them.
        full_para_text = paragraph.text
        matches = re.findall(r"\{\{.*?\}\}", full_para_text)
        if not matches:
            continue

        left_val   = (shape.left   or 0) * scale_x + offset_x
        top_val    = (shape.top    or 0) * scale_y + offset_y
        width_val  = (shape.width  or 0) * scale_x
        height_val = (shape.height or 0) * scale_y

        left_pct   = (left_val   / slide_width)  * 100
        top_pct    = (top_val    / slide_height) * 100
        width_pct  = (width_val  / slide_width)  * 100 if width_val else 10
        height_pct = (height_val / slide_height) * 100 if height_val else 10

        font_size_pt   = None
        font_color_hex = None
        is_bold        = False
        is_italic      = False

        # Scan all runs for font metadata — even split runs carry valid styling.
        for run in paragraph.runs:
            if run.font:
                if not font_size_pt and run.font.size:
                    font_size_pt = run.font.size.pt
                if not font_color_hex and hasattr(run.font, "color"):
                    try:
                        if hasattr(run.font.color, "rgb") and run.font.color.rgb:
                            font_color_hex = f"#{str(run.font.color.rgb)}"
                    except Exception:
                        pass
                if run.font.bold:
                    is_bold = True
                if run.font.italic:
                    is_italic = True

        if not font_size_pt and paragraph.font and paragraph.font.size:
            font_size_pt = paragraph.font.size.pt
        if not font_size_pt:
            font_size_pt = 18
        if not font_color_hex:
            font_color_hex = "#1e293b"

        align_val = "CENTER"
        raw_align = paragraph.alignment
        if raw_align is None:
            try:
                pPr = paragraph._p.find('{http://schemas.openxmlformats.org/drawingml/2006/main}pPr')
                if pPr is not None:
                    algn = pPr.get('algn')
                    if algn == 'r':
                        align_val = "RIGHT"
                    elif algn == 'l':
                        align_val = "LEFT"
                    else:
                        align_val = "CENTER"
            except Exception:
                align_val = "CENTER"
        else:
            align_str = str(raw_align)
            if "CENTER" in align_str:
                align_val = "CENTER"
            elif "RIGHT" in align_str:
                align_val = "RIGHT"
            else:
                align_val = "LEFT"

        for match in matches:
            if match not in placeholders_map:
                placeholders_map[match] = {
                    "left":     left_pct,
                    "top":      top_pct,
                    "width":    width_pct,
                    "height":   height_pct,
                    "fontSize": font_size_pt,
                    "align":    align_val,
                    "color":    font_color_hex,
                    "isBold":   is_bold,
                    "isItalic": is_italic,
                }

def main():
    try:
        input_data = json.loads(sys.stdin.read())
        url = input_data.get("templateUrl")
        if not url:
            raise ValueError("Template URL is missing")

        # Download file
        response = requests.get(url, timeout=30)
        response.raise_for_status()

        # ── Detect file type ──────────────────────────────────────────
        url_lower = url.lower().split('?')[0]  # strip query params
        is_svg = url_lower.endswith('.svg') or b'<svg' in response.content[:500]

        if is_svg:
            sys.stderr.write("Detected SVG file — using SVG parser\n")
            svg_text = response.content.decode('utf-8', errors='replace')
            placeholders_map, svg_width, svg_height = extract_svg_placeholders(svg_text)

            print(json.dumps({
                "success": True,
                "file_type": "svg",
                "placeholders": list(placeholders_map.keys()),
                "placeholder_mappings": placeholders_map,
                "slide_width": svg_width,
                "slide_height": svg_height
            }))

        else:
            # ── PPTX path ─────────────────────────────────────────────
            # Validate PPTX magic bytes
            if len(response.content) < 100 or response.content[:2] != b'PK':
                raise ValueError(f"File is not a valid PPTX (bad magic bytes)")

            with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp_in:
                tmp_in.write(response.content)
                tmp_in_path = tmp_in.name

            try:
                from pptx import Presentation
                prs = Presentation(tmp_in_path)
                placeholders_map = {}
                slide_width = prs.slide_width
                slide_height = prs.slide_height

                for slide_idx, slide in enumerate(prs.slides):
                    slide_placeholders_before = set(placeholders_map.keys())
                    for shape in slide.shapes:
                        extract_placeholders(shape, slide_width, slide_height, placeholders_map)
                    # Tag any newly discovered placeholders with this slide's index
                    for ph, meta in placeholders_map.items():
                        if ph not in slide_placeholders_before and 'slide_index' not in meta:
                            meta['slide_index'] = slide_idx

                sys.stderr.write(f"Detected placeholders: {list(placeholders_map.keys())}\n")

                print(json.dumps({
                    "success": True,
                    "file_type": "pptx",
                    "placeholders": list(placeholders_map.keys()),
                    "placeholder_mappings": placeholders_map,
                    "slide_width": slide_width,
                    "slide_height": slide_height
                }))
            finally:
                if os.path.exists(tmp_in_path):
                    os.remove(tmp_in_path)

    except Exception as e:
        print(json.dumps({"success": False, "error": str(e)}))
        sys.exit(1)

if __name__ == "__main__":
    main()